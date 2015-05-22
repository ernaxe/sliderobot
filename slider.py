'''
Created on 12 Sep 2014

@author: hpais
'''
import os
from pptx import Presentation
from pptx import util
from pptx.enum.shapes import MSO_SHAPE

SLIDE_WIDTH = util.Cm(25.4)
SLIDE_HEIGHT = util.Cm(19.05)
SLIDE_LAYOUT_BLANK = 6

def image_folder2pptx(fd,ppname,nrow,ncol,picwidth,picheight,\
                      resize=False, reposition = True):
    if resize:
        step_x = SLIDE_WIDTH/ncol
        step_y = SLIDE_HEIGHT/nrow
    else:
        step_x = util.Px(picwidth)
        step_y = util.Px(picheight)
        #print("test")
        #print step_x, step_y, picwidth, picheight
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[SLIDE_LAYOUT_BLANK])
    pos_x = 0
    pos_y = 0
    k = 1
    for fn in os.listdir(fd):
        print fn, k
        slide.shapes.add_picture(fd+fn,pos_x,pos_y,step_x,step_y)
        if reposition:
            pos_x += step_x
            if k % ncol == 0:
                pos_x = 0
                pos_y += step_y
                if k % (nrow*ncol) == 0:
                    pos_y = 0
                    slide = prs.slides.add_slide(
                        prs.slide_layouts[SLIDE_LAYOUT_BLANK])
        k += 1
        #if  and (not resize):
        #    slide = prs.slides.add_slide(prs.slide_layouts[SLIDE_LAYOUT_BLANK])
    prs.save(ppname)
